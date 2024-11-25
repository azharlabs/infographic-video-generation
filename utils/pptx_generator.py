import os
import logging
from io import BytesIO
from pptx import Presentation
from .animation_generator import generate_slide_animation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Configure logging
logger = logging.getLogger(__name__)


def calculate_textbox_height(text,
                             lines,
                             min_height=Inches(1.0),
                             max_height=Inches(4.5)):
    average_line_height = Pt(18)  # Estimated average height per line
    estimated_height = average_line_height * len(lines)
    estimated_height = max(min_height, min(max_height, estimated_height))
    return estimated_height


def adjust_text_box(text_frame,
                    max_size=40,
                    min_size=10,
                    color=None,
                    alignment=PP_ALIGN.LEFT):
    """Adjust text size to fit in text box with proper styling"""
    if not hasattr(text_frame, 'paragraphs'):
        logger.warning("Text frame does not have paragraphs attribute")
        return

    for paragraph in text_frame.paragraphs:
        paragraph.alignment = alignment
        if not hasattr(paragraph, 'runs'):
            continue

        for run in paragraph.runs:
            if not hasattr(run, 'font'):
                continue

            font = run.font
            # Dynamic sizing based on content length
            content_length = len(paragraph.text) if hasattr(paragraph,
                                                            'text') else 0
            font.size = Pt(
                min(max_size, max(min_size, 40 - content_length // 20)))

            if color and hasattr(font, 'color'):
                font.color.rgb = color


def add_animation_to_slide(slide, animation_path, position):
    """Add animation GIF to slide with proper positioning and validation"""
    if not animation_path or not os.path.exists(animation_path):
        logger.error(
            f"Animation file not found or invalid path: {animation_path}")
        return False

    try:
        # Verify GIF format and readability
        print(animation_path)
        with open(animation_path, 'rb') as f:
            header = f.read(6)
            if not header.startswith(b'GIF8'):
                logger.error(f"Invalid GIF format for {animation_path}")
                return False

        # Get file size and check permissions
        file_size = os.path.getsize(animation_path)
        if file_size < 100:
            logger.error(
                f"Animation file too small: {animation_path} ({file_size} bytes)"
            )
            return False

        if not os.access(animation_path, os.R_OK):
            logger.error(f"No read permission for {animation_path}")
            return False

        # Get presentation dimensions from parent presentation
        prs = slide.part.package.presentation
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Calculate dimensions based on slide size
        left = position.get('left', Inches(7))  # Positioned on the right side
        top = position.get('top', Inches(2))  # Centered vertically
        width = position.get('width',
                             min(Inches(5),
                                 slide_width * 0.4))  # Max 40% of slide width
        height = position.get('height', min(Inches(4), slide_height *
                                            0.6))  # Max 60% of slide height

        with open(animation_path, 'rb') as gif:
            shape = slide.shapes.add_picture(gif, left, top, width, height)

        logger.info(f"Successfully added animation to slide: {animation_path}")
        return True

    except Exception as e:
        logger.error(f"Error adding animation to slide: {str(e)}")
        return False


def create_presentation(
        content: str,
        template_name: str = "modern",
        animations: list[BytesIO] | None = None) -> Presentation:
    """Generate PPTX file from content with specific formatting and animations"""
    logger.info(f"Creating presentation with template: {template_name}")

    # Generate animations for each section if not provided
    if not animations:
        animations = []
        sections = content
        for i, section in enumerate(sections):
            animation_path = os.path.join(
                os.path.dirname(os.path.dirname(__file__)), 'static',
                'animations', f'slide_{i}_animation.gif')
            os.makedirs(os.path.dirname(animation_path), exist_ok=True)

            # Create animation based on section content
            animation_gif = generate_slide_animation(section)
            # Convert gif data to BytesIO if it's not already
            if isinstance(animation_gif, str):
                with open(animation_gif, 'rb') as file_data:
                    animation_gif = BytesIO(file_data.read())
            animations.append(animation_gif)
    """Generate PPTX file from content with enhanced formatting and animations"""
    prs = Presentation()

    # Set slide dimensions for widescreen format (16:9)
    slide_width = Inches(13.33)
    slide_height = Inches(7.5)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    # Define font styles based on template
    font_styles = {
        'title': {
            'name': 'Calibri',
            'size': Pt(24),
            'bold': True
        },
        'body': {
            'name': 'Calibri',
            'size': Pt(20),
            'bold': False
        }
    }

    # Template configurations
    templates = {
        'modern': {
            'background': RGBColor(240, 240, 240),
            'title_color': RGBColor(51, 51, 51),
            'accent_color': RGBColor(108, 99, 255)
        },
        'professional': {
            'background': RGBColor(255, 255, 255),
            'title_color': RGBColor(0, 0, 0),
            'accent_color': RGBColor(0, 102, 204)
        },
        'minimal': {
            'background': RGBColor(255, 255, 255),
            'title_color': RGBColor(51, 51, 51),
            'accent_color': RGBColor(0, 0, 0)
        },
        'gradient': {
            'background': RGBColor(42, 42, 42),
            'title_color': RGBColor(255, 255, 255),
            'accent_color': RGBColor(108, 99, 255)
        },
        'corporate': {
            'background': RGBColor(31, 58, 99),
            'title_color': RGBColor(255, 255, 255),
            'accent_color': RGBColor(255, 255, 255)
        },
        'creative': {
            'background': RGBColor(42, 42, 42),
            'title_color': RGBColor(255, 107, 107),
            'accent_color': RGBColor(78, 205, 196)
        },
        'dynamic': {
            'background': RGBColor(45, 45, 45),
            'title_color': RGBColor(255, 255, 255),
            'accent_color': RGBColor(255, 126, 0)
        },
        'clean': {
            'background': RGBColor(248, 249, 250),
            'title_color': RGBColor(33, 37, 41),
            'accent_color': RGBColor(13, 110, 253)
        },
        'dark': {
            'background': RGBColor(18, 18, 18),
            'title_color': RGBColor(255, 255, 255),
            'accent_color': RGBColor(130, 177, 255)
        },
        'tech': {
            'background': RGBColor(22, 28, 36),
            'title_color': RGBColor(0, 255, 255),
            'accent_color': RGBColor(64, 196, 255)
        }
    }
    template_config = templates.get(template_name, templates['modern'])
    bg_color = template_config['background']

    # Add title slide with background
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    # Configure title slide
    title_shape = slide.shapes.title
    subtitle_shape = None
    if len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.top = Inches(6)

    # Extract title from first section if available
    sections = content.split('\n\n')
    main_title = sections[0].split('\n')[0] if sections else "Presentation"

    if title_shape:
        title_shape.text = main_title
        # Style title slide
        title_tf = title_shape.text_frame
        title_tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Apply template styles to title
        for paragraph in title_tf.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                font = run.font
                font.name = font_styles['title']['name']
                font.size = font_styles['title']['size']
                font.bold = font_styles['title']['bold']
                font.color.rgb = template_config['title_color']

                # Measure the length of text and compare with shape width
                if title_shape:
                    title_width = title_shape.width
                    text_length = Pt(len(title_shape.text) *
                                     9)  # Roughly 9pt per char
                    if title_width < text_length:
                        title_shape.text = title_shape.text + "\n"

    if subtitle_shape:
        try:
            subtitle_shape.text = "Project by team Turing-1950"
            # Apply template styles to subtitle
            subtitle_tf = subtitle_shape.text_frame
            for paragraph in subtitle_tf.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    font = run.font
                    font.name = font_styles['body']['name']
                    font.size = Pt(24)  # Smaller than title
                    font.color.rgb = template_config['accent_color']
        except AttributeError:
            logger.warning("Subtitle shape does not support text attribute")

    # Add content slides
    animations = animations or []
    for i, section in enumerate(sections):
        content_slide_layout = prs.slide_layouts[
            1]  # Using layout with title and content
        slide = prs.slides.add_slide(content_slide_layout)

        # Set slide dimensions for 16:9
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Split section into lines and format title
        lines = section.strip().split('\n')
        title_shape = slide.shapes.title
        if title_shape and lines:
            title_shape.text = lines[0]
            title_shape.top = Inches(0.8)
            title_shape.left = Inches(1.5)
            title_shape.width = Inches(10.3)

            title_frame = title_shape.text_frame
            title_frame.word_wrap = True  # Enable word wrap to prevent overflow

            # Set title paragraph properties
            paragraph = title_frame.paragraphs[0]
            paragraph.font.size = Pt(36)
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.name = 'Arial'

        # Add image on left side
        if i < len(animations) and animations[i]:
            left_image = Inches(0.75)
            top_image = Inches(2.0)
            width_image = Inches(5.5)
            height_image = Inches(4.8)

            slide.shapes.add_picture(animations[i], left_image, top_image,
                                     width_image, height_image)

        # Add text box for bullet points on right side
        left_text = Inches(6.75)
        top_text = Inches(2.0)
        width_text = Inches(6)
        height_text = Inches(4.8)

        text_box = slide.shapes.add_textbox(left_text, top_text, width_text,
                                            height_text)

        # Format bullet points
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        # Add bullet points from remaining lines
        if len(lines) > 1:
            bullet_points = lines[1:]
            for idx, point in enumerate(bullet_points):
                p = text_frame.add_paragraph()
                p.text = point.strip()
                p.font.size = Pt(24)
                p.font.name = 'Arial'
                p.level = 0
                if i > 0:
                    p.space_before = Pt(12)
        # Add template-specific design elements
        if template_name == 'modern':
            # Add accent bar on the left
            accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                 Inches(0), Inches(0),
                                                 Inches(0.2), slide_height)
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = template_config['accent_color']
            accent_line.line.fill.background()

        elif template_name == 'gradient':
            # Add gradient overlay
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                             Inches(0), slide_width,
                                             slide_height)
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = template_config['accent_color']
            overlay.fill.transparency = 0.85

        elif template_name == 'corporate':
            # Add header bar
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                            Inches(0), slide_width, Inches(1))
            header.fill.solid()
            header.fill.fore_color.rgb = template_config['accent_color']
            header.line.fill.background()

        elif template_name == 'tech':
            # Add tech pattern
            for i in range(3):
                line = slide.shapes.add_shape(MSO_SHAPE.LINE, Inches(0.5 + i),
                                              Inches(6), Inches(2), Inches(0))
                line.line.color.rgb = template_config['accent_color']
                line.line.transparency = 0.7
        # Add background to slide (after content to ensure proper z-order)
        background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0),
                                            Inches(0), slide_width,
                                            slide_height)
        background.fill.solid()
        background.fill.fore_color.rgb = bg_color
        background.line.fill.background()
        background.z_order = 0  # Put background at the bottom

        # Add accent line or shape based on template
        if template_name in ['modern', 'corporate', 'tech']:
            accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                 Inches(0), Inches(0.5),
                                                 Inches(0.2), Inches(6.5))
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = template_config['accent_color']
            accent_line.line.fill.background()

        # First create all background and design elements
        lines = section.split('\n')

        # Generate and add animation
        animation_path = generate_slide_animation(section)
        if animation_path and os.path.exists(animation_path):
            # Add visualization to slide - positioned on the left side
            left_image = Inches(0.75)
            top_image = Inches(2.0)
            width_image = Inches(5.5)
            height_image = Inches(4.8)

            try:
                # Add animation using file path
                if os.path.exists(animation_path):
                    slide.shapes.add_picture(animation_path, left_image,
                                             top_image, width_image,
                                             height_image)
                    logger.info(
                        f"Successfully added animation to slide: {animation_path}"
                    )
                else:
                    logger.warning(
                        f"Animation file not found: {animation_path}")
            except Exception as e:
                logger.error(f"Error adding animation to slide: {str(e)}")

        # Store references to shapes we'll need to modify later
        title_placeholder = slide.shapes.title
        body_placeholder = slide.placeholders[1] if len(
            slide.placeholders) > 1 else None

        # Remove existing title and body placeholders temporarily
        if title_placeholder:
            sp = title_placeholder._element
            sp.getparent().remove(sp)
        if body_placeholder:
            sp = body_placeholder._element
            sp.getparent().remove(sp)

        # Create text shapes last to ensure they're on top
        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5),
                                               Inches(9), Inches(1.2))
        title_shape.z_order = 2  # Ensure title is on top

        # Determine the lines of text
        text_lines = lines[1:] if len(lines) > 1 else []
        combined_text = '\n'.join(text_lines)
        # Calculate the dynamic height
        body_textbox_height = calculate_textbox_height(combined_text,
                                                       text_lines)
        # Add the textbox with dynamic height
        body_shape = slide.shapes.add_textbox(Inches(6.75), Inches(1.8),
                                              Inches(6), Inches(6))
        body_shape.z_order = 2

        # Set title if shape exists and has required attributes
        if title_shape and lines and hasattr(title_shape, 'text_frame'):
            try:
                title_tf = title_shape.text_frame
                title_tf.clear()  # Clear existing text
                title_tf.margin_bottom = Inches(0.1)
                title_tf.margin_left = Inches(0.5)

                p = title_tf.paragraphs[0]
                p.text = lines[0]
                p.alignment = PP_ALIGN.LEFT

                # Word wrap logic if text overflows shape
                content_length = len(p.text)
                title_width = title_shape.width if hasattr(
                    title_shape, 'width') else 0
                if content_length > title_width / 8:  # Assuming average character width
                    wrapped_text = '\n'.join(
                        p.text[i:i + int(title_width / 8)]
                        for i in range(0, content_length, int(title_width /
                                                              8)))
                    p.text = wrapped_text

                for run in p.runs:
                    if hasattr(run, 'font'):
                        font = run.font
                        font.name = font_styles['title']['name']
                        font.size = Pt(28)  # Slightly smaller than title slide
                        font.bold = font_styles['title']['bold']
                        font.color.rgb = template_config['title_color']
            except Exception as e:
                logger.warning(f"Error setting title text: {str(e)}")

        # Set content if body shape exists and has text frame
        if body_shape and len(lines) > 1:
            try:
                if not hasattr(body_shape, 'text_frame'):
                    logger.warning(
                        "Body shape does not have text_frame attribute")
                    continue
                content_lines = lines[1:]
                tf = body_shape.text_frame
                tf.clear()  # Clear existing text
                # Set proper margins
                tf.margin_left = Inches(0.5)
                tf.margin_right = Inches(0.5)
                tf.margin_top = Inches(0.2)
                # Add content with proper styling
                p = tf.paragraphs[0]
                for idx, line in enumerate(content_lines):
                    if idx > 0:
                        p.add_line_break()  # Add line break between lines
                    run = p.add_run()
                    run.text = line
                    font = run.font
                    font.name = font_styles['body']['name']
                    font.size = Pt(24)
                    font.color.rgb = template_config[
                        'title_color']  # Use title color for better contrast
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(12)  # Space between paragraphs

            except AttributeError:
                logger.warning("Body shape does not support text frame")

    return prs
