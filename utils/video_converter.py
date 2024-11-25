import os
import gc
import logging
from pathlib import Path
import mimetypes
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import ColorFormat, RGBColor
from pptx.enum.dml import MSO_FILL_TYPE, MSO_THEME_COLOR_INDEX
from moviepy.editor import *
from PIL import Image, ImageDraw, ImageFont, ImageSequence
import numpy as np
import subprocess
from io import BytesIO

# Configure logging with more detailed format
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)


def get_font_color(run):
    """Helper function to safely get font color"""
    try:
        if hasattr(run, 'font') and hasattr(run.font, 'color'):
            if run.font.color.type == MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
                if hasattr(run.font.color,
                           'rgb') and run.font.color.rgb is not None:
                    return (run.font.color.rgb[0], run.font.color.rgb[1],
                            run.font.color.rgb[2], 255)
    except AttributeError:
        pass
    # Default to black if no color is specified
    return (0, 0, 0, 255)


def process_shape_text(shape):
    try:
        text = ""
        if hasattr(shape, 'text_frame') and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text
                text += '\n'
        elif hasattr(shape, 'text') and shape.text:
            text = shape.text
        return text.strip()
    except Exception as e:
        logger.error(f"Error processing shape text: {str(e)}")
        return ""


def verify_file_type(file_path: str, expected_type: str) -> bool:
    """Verify if file is of expected type using mimetypes"""
    if not os.path.exists(file_path):
        logger.error(f"File does not exist: {file_path}")
        return False

    try:
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            logger.warning(f"Could not determine mime type for: {file_path}")
            return False

        logger.debug(f"File {file_path} type: {mime_type}")
        return expected_type in mime_type.lower()
    except Exception as e:
        logger.error(f"Error verifying file type: {str(e)}")
        return False


def check_imagemagick():
    """Check if ImageMagick is available and properly configured"""
    try:
        # Check ImageMagick installation
        result = subprocess.run(['convert', '-version'],
                                check=True,
                                capture_output=True,
                                text=True)
        logger.info(f"ImageMagick version: {result.stdout.split()[2]}")
        return True
    except (subprocess.CalledProcessError, FileNotFoundError) as e:
        logger.error(f"ImageMagick check failed: {str(e)}")
        return False


def apply_transition_effect(clip, transition_type='fade', duration=0.5):
    """Apply transition effect to a clip"""
    try:
        if transition_type == 'fade':
            # Apply fade in and fade out
            clip = clip.fadein(duration).fadeout(duration)
        elif transition_type == 'slide_left':
            # Slide from left transition
            w, h = clip.size
            mask = ColorClip(clip.size, col=(0, 0, 0))
            mask = mask.set_duration(clip.duration)
            sliding = clip.set_position(
                lambda t: (max(0, int(w * (t / duration - 1))), 0))
            clip = CompositeVideoClip([mask, sliding])
        elif transition_type == 'slide_right':
            # Slide from right transition
            w, h = clip.size
            mask = ColorClip(clip.size, col=(0, 0, 0))
            mask = mask.set_duration(clip.duration)
            sliding = clip.set_position(
                lambda t: (min(0, int(w * (1 - t / duration))), 0))
            clip = CompositeVideoClip([mask, sliding])
        elif transition_type == 'zoom':
            # Zoom in transition
            clip = clip.resize(lambda t: 1 + 0.1 * t / duration)

        return clip
    except Exception as e:
        logger.error(f"Error applying transition effect: {str(e)}")
        return clip


def create_slide_clip(slide, width, height, duration, slide_width_emus,
                      slide_height_emus):
    try:
        # Create base frame with slide background
        background_image = Image.new('RGB', (width, height), (255, 255, 255))
        background_fill = None

        # Get background fill from slide background
        if hasattr(slide, 'background') and slide.background:
            if hasattr(slide.background, 'fill'):
                background_fill = slide.background.fill

        # If no slide background, try to get background from slide layout
        if not background_fill and hasattr(
                slide, 'slide_layout') and slide.slide_layout:
            if hasattr(slide.slide_layout,
                       'background') and slide.slide_layout.background:
                if hasattr(slide.slide_layout.background, 'fill'):
                    background_fill = slide.slide_layout.background.fill

        # If still no background, try to get from master slide
        if not background_fill and hasattr(
                slide.slide_layout,
                'slide_master') and slide.slide_layout.slide_master:
            master = slide.slide_layout.slide_master
            if hasattr(master, 'background') and master.background:
                if hasattr(master.background, 'fill'):
                    background_fill = master.background.fill

        if background_fill:
            logger.debug(f"Background fill type: {background_fill.type}")

            if background_fill.type == MSO_FILL_TYPE.SOLID:
                # Solid color background
                fore_color = background_fill.fore_color
                if hasattr(fore_color, 'rgb'):
                    color = fore_color.rgb
                    background_color = (color[0], color[1], color[2])
                    background_image = Image.new('RGB', (width, height),
                                                 background_color)
                    logger.debug(
                        f"Applied solid background color: {background_color}")

            elif background_fill.type == MSO_FILL_TYPE.GRADIENT:
                # Handle gradient background
                logger.debug("Processing gradient background")
                background_image = create_gradient_background(
                    background_fill, width, height)

            elif background_fill.type == MSO_FILL_TYPE.PICTURE:
                try:
                    # Picture background
                    image_blob = background_fill.picture.image.blob
                    bg_image_stream = BytesIO(image_blob)
                    bg_pil_image = Image.open(bg_image_stream)
                    bg_pil_image = bg_pil_image.resize(
                        (width, height), Image.Resampling.LANCZOS)
                    background_image = bg_pil_image.convert('RGB')
                    logger.debug("Applied picture background")
                except Exception as e:
                    logger.error(
                        f"Error processing picture background: {str(e)}")

            else:
                # Default white background for unsupported fill types
                logger.debug(
                    f"Unsupported or unhandled fill type. Using default white background."
                )
                background_image = Image.new('RGB', (width, height),
                                             (255, 255, 255))

        # Convert background image to numpy array
        frame = np.array(background_image)
        background_clip = ImageClip(frame).set_duration(duration)
        clips = [background_clip]

        # Process all elements in the slide
        for shape in slide.shapes:

            logger.debug(
                f"Processing shape: name={shape.name}, type={shape.shape_type}"
            )
            try:
                if hasattr(shape, 'fill'):
                    fill = shape.fill
                    if fill.type == MSO_FILL_TYPE.SOLID:
                        # Create shape mask
                        shape_image = Image.new('RGBA', (width, height),
                                                (0, 0, 0, 0))
                        draw = ImageDraw.Draw(shape_image)

                        # Get shape dimensions
                        left = shape.left if hasattr(shape, 'left') else 0
                        top = shape.top if hasattr(shape, 'top') else 0
                        shape_width = shape.width if hasattr(
                            shape, 'width') else width
                        shape_height = shape.height if hasattr(
                            shape, 'height') else height

                        # Convert EMUs to pixels
                        x = int((left / slide_width_emus) * width)
                        y = int((top / slide_height_emus) * height)
                        w = int((shape_width / slide_width_emus) * width)
                        h = int((shape_height / slide_height_emus) * height)

                        # Draw shape
                        color = fill.fore_color.rgb
                        shape_color = (color[0], color[1], color[2], 255)
                        draw.rectangle([x, y, x + w, y + h], fill=shape_color)

                        # Convert to numpy array and create clip
                        shape_array = np.array(shape_image)
                        shape_clip = ImageClip(shape_array).set_duration(
                            duration)
                        clips.append(shape_clip)

                # Process shapes that have a text frame
                if hasattr(shape,
                           'text_frame') and shape.text_frame is not None:
                    logger.debug("Processing text frame")
                    text_image = Image.new('RGBA', (width, height),
                                           (0, 0, 0, 0))
                    draw = ImageDraw.Draw(text_image)
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text.strip()
                            if not text:
                                continue
                            # Get font size and color
                            font_size = run.font.size.pt if run.font.size else 24
                            font_color = get_font_color(run)

                            # Use a system font
                            try:
                                font = ImageFont.truetype(
                                    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                                    int(font_size))
                            except Exception as e:
                                logger.warning(
                                    f"Font not found. Using default font. Error: {e}"
                                )
                                font = ImageFont.load_default()

                            # Get shape position
                            left = shape.left if hasattr(shape, 'left') else 0
                            top = shape.top if hasattr(shape, 'top') else 0

                            # Convert EMUs to pixels
                            x = int((left / slide_width_emus) * width)
                            y = int((top / slide_height_emus) * height)

                            # Draw text
                            draw.text((x, y), text, font=font, fill=font_color)
                            logger.debug(
                                f"Drew text: '{text}' at ({x}, {y}) with color {font_color}"
                            )

                    # Convert to numpy array
                    text_image_np = np.array(text_image)

                    # Create text clip
                    text_clip = ImageClip(text_image_np).set_duration(duration)
                    clips.append(text_clip)

                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if hasattr(shape, 'image'):
                        # Handle images
                        image_stream = BytesIO(shape.image.blob)
                        pil_image = Image.open(image_stream)

                        # Check if image is a GIF
                        if pil_image.format == 'GIF' and getattr(
                                pil_image, 'is_animated', False):
                            logger.debug("Processing animated GIF")
                            # Get image position and size
                            left = shape.left if hasattr(shape, 'left') else 0
                            top = shape.top if hasattr(shape, 'top') else 0
                            shape_width = shape.width if hasattr(
                                shape, 'width') else slide_width_emus // 2
                            shape_height = shape.height if hasattr(
                                shape, 'height') else slide_height_emus // 2

                            x = int((left / slide_width_emus) * width)
                            y = int((top / slide_height_emus) * height)
                            w = int((shape_width / slide_width_emus) * width)
                            h = int(
                                (shape_height / slide_height_emus) * height)

                            # Extract frames from GIF
                            frames = []
                            durations = []
                            for frame in ImageSequence.Iterator(pil_image):
                                frame = frame.convert('RGBA')
                                frame = frame.resize((w, h),
                                                     Image.Resampling.LANCZOS)
                                frame_np = np.array(frame)
                                frames.append(frame_np)
                                durations.append(
                                    frame.info.get('duration', 100) /
                                    1000.0)  # Duration in seconds

                            # Create an ImageSequenceClip
                            gif_clip = ImageSequenceClip(frames,
                                                         durations=durations)
                            gif_clip = gif_clip.set_position(
                                (x, y)).set_duration(duration)
                            clips.append(gif_clip)
                        else:
                            logger.debug("Processing static image")
                            # Handle static images
                            if pil_image.mode != 'RGBA':
                                pil_image = pil_image.convert('RGBA')

                            # Get image position and size
                            left = shape.left if hasattr(shape, 'left') else 0
                            top = shape.top if hasattr(shape, 'top') else 0
                            shape_width = shape.width if hasattr(
                                shape, 'width') else slide_width_emus // 2
                            shape_height = shape.height if hasattr(
                                shape, 'height') else slide_height_emus // 2

                            # Convert EMUs to pixels
                            x = int((left / slide_width_emus) * width)
                            y = int((top / slide_height_emus) * height)
                            w = int((shape_width / slide_width_emus) * width)
                            h = int(
                                (shape_height / slide_height_emus) * height)

                            # Resize image
                            pil_image = pil_image.resize(
                                (w, h), Image.Resampling.LANCZOS)

                            # Create new image with correct position
                            full_image = Image.new('RGBA', (width, height),
                                                   (0, 0, 0, 0))
                            full_image.paste(pil_image, (x, y), pil_image)

                            # Convert to numpy array and create clip
                            image_array = np.array(full_image)
                            image_clip = ImageClip(image_array).set_duration(
                                duration)
                            clips.append(image_clip)

                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    if hasattr(shape, 'table'):
                        # Handle tables
                        table_image = Image.new('RGBA', (width, height),
                                                (0, 0, 0, 0))
                        draw = ImageDraw.Draw(table_image)

                        try:
                            font = ImageFont.truetype(
                                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                                24)
                        except Exception as e:
                            logger.warning(
                                f"Font not found. Using default font. Error: {e}"
                            )
                            font = ImageFont.load_default()

                        # Draw table
                        num_rows = len(shape.table.rows)
                        num_cols = len(shape.table.columns)
                        cell_height = height // num_rows
                        cell_width = width // num_cols

                        for row_idx, row in enumerate(shape.table.rows):
                            for col_idx, cell in enumerate(row.cells):
                                x = col_idx * cell_width
                                y = row_idx * cell_height

                                # Draw cell background
                                draw.rectangle(
                                    [x, y, x + cell_width, y + cell_height],
                                    fill=(255, 255, 255,
                                          255),  # Opaque background
                                    outline=(0, 0, 0, 255))

                                # Draw cell text
                                text = cell.text.strip()
                                if text:
                                    draw.text((x + 5, y + 5),
                                              text,
                                              font=font,
                                              fill=(0, 0, 0, 255))

                        table_array = np.array(table_image)
                        table_clip = ImageClip(table_array).set_duration(
                            duration)
                        clips.append(table_clip)

            except Exception as shape_error:
                logger.error(f"Error processing shape: {str(shape_error)}",
                             exc_info=True)
                continue

        # Combine all clips
        if len(clips) > 1:
            composite_clip = CompositeVideoClip(clips, size=(width, height))
            return composite_clip
        return clips[0]

    except Exception as e:
        logger.error(f"Error creating slide clip: {str(e)}", exc_info=True)
        return background_clip


def create_pattern_background(fore_color, back_color, width, height):
    """Create a pattern background"""
    try:
        # Create a simple pattern (checkerboard as an example)
        pattern_size = 20
        background = Image.new('RGB', (width, height), back_color)
        draw = ImageDraw.Draw(background)

        for x in range(0, width, pattern_size * 2):
            for y in range(0, height, pattern_size * 2):
                draw.rectangle([x, y, x + pattern_size, y + pattern_size],
                               fill=fore_color)
                draw.rectangle([
                    x + pattern_size, y + pattern_size, x + pattern_size * 2,
                    y + pattern_size * 2
                ],
                               fill=fore_color)

        return background
    except Exception as e:
        logger.error(f"Error creating pattern background: {str(e)}")
        return Image.new('RGB', (width, height), (255, 255, 255))


def convert_pptx_to_video(pptx_path: str, output_path: str) -> None:
    try:
        # Load presentation
        prs = Presentation(pptx_path)
        logger.info(f"Loaded presentation with {len(prs.slides)} slides")

        # Video settings
        width = 1920
        height = 1080
        fps = 24

        # Get slide dimensions
        slide_width_emus = prs.slide_width
        slide_height_emus = prs.slide_height

        transitions = ['fade', 'slide_left', 'slide_right', 'zoom']

        # Process each slide
        clips = []
        for i, slide in enumerate(prs.slides, 1):
            try:
                logger.info(f"Processing slide {i}")

                # Set slide duration to 2 seconds
                duration = 2
                transition_duration = 0.5  # Duration for transitions
                total_duration = duration + transition_duration

                # Create slide clip with slide dimensions
                clip = create_slide_clip(slide, width, height, total_duration,
                                         slide_width_emus, slide_height_emus)

                # Add transition effect
                if clip is not None:
                    transition_type = transitions[
                        i % len(transitions)]  # Cycle through transitions
                    clip = apply_transition_effect(
                        clip,
                        transition_type=transition_type,
                        duration=transition_duration)
                    clips.append(clip)

                    # Clean up memory after each slide
                    gc.collect()
                    logger.info(f"Done slide {i}")

            except Exception as slide_error:
                logger.error(f"Error processing slide {i}: {str(slide_error)}",
                             exc_info=True)
                # Add blank slide as fallback
                frame = np.ones((height, width, 3), dtype=np.uint8) * 255
                clips.append(
                    ImageClip(frame).set_duration(transition_duration))

        if not clips:
            raise ValueError("No valid slides were processed")

        # Concatenate all slides
        final_clip = concatenate_videoclips(clips, method="compose")

        # Write video file
        final_clip.write_videofile(
            output_path,
            fps=fps,
            codec='libx264',
            audio=False,
            preset='ultrafast',
            threads=4,
            logger=None,
            bitrate='2000k',  # Added bitrate controlj
            ffmpeg_params=[
                '-tune',
                'fastdecode',
                '-movflags',
                '+faststart',
                '-bf',
                '0'  # Disable B-frames for faster encoding
            ])

        # Clean up
        final_clip.close()
        for clip in clips:
            clip.close()
        gc.collect()

    except Exception as e:
        logger.error(f"Error in conversion: {str(e)}", exc_info=True)
        raise RuntimeError(f"Failed to convert PPTX to video: {str(e)}")


def convert_to_video(pptx_path: str, output_path: str) -> bool:
    """Main conversion function"""
    logger.info(f"Starting video conversion for PPTX: {pptx_path}")

    try:
        # Check ImageMagick
        if not check_imagemagick():
            raise RuntimeError("ImageMagick is not properly configured")

        # Verify input file
        if not verify_file_type(pptx_path, 'officedocument'):
            raise ValueError("Invalid PPTX file format")

        # Ensure paths are absolute
        pptx_path = os.path.abspath(pptx_path)
        output_path = os.path.abspath(output_path)

        # Create output directory if needed
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        # Convert PPTX to video
        convert_pptx_to_video(pptx_path, output_path)

        # Verify output
        if not os.path.exists(output_path):
            raise FileNotFoundError(f"Video was not created at {output_path}")

        video_size = os.path.getsize(output_path)
        if video_size < 1000:
            raise ValueError(
                f"Generated video is too small ({video_size} bytes)")

        logger.info(
            f"Successfully created video: {output_path} ({video_size} bytes)")
        return True

    except Exception as e:
        logger.error(f"Conversion failed: {str(e)}", exc_info=True)
        if os.path.exists(output_path):
            try:
                os.remove(output_path)
            except Exception:
                pass
        raise
